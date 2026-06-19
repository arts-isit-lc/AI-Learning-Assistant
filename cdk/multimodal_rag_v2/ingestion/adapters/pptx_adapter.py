"""PPTX adapter — extracts text, images, and tables from PowerPoint files."""

from __future__ import annotations

from io import BytesIO
from typing import TYPE_CHECKING

from aws_lambda_powertools import Logger
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ...models.data_models import ElementType, Provenance, RawElement
from ..base_adapter import BaseAdapter

if TYPE_CHECKING:
    from ...models.data_models import FileMetadata

logger = Logger(service="multimodal-rag-ingestion")


class PptxAdapter(BaseAdapter):
    """Extract text, images, and tables from PPTX files using python-pptx."""

    def extract(
        self, file_content: bytes, file_metadata: FileMetadata
    ) -> list[RawElement]:
        """Extract elements from all slides in a PPTX file.

        Args:
            file_content: Raw bytes of the PPTX file.
            file_metadata: Metadata about the file.

        Returns:
            List of RawElement instances with slide-level provenance.
        """
        presentation = Presentation(BytesIO(file_content))
        elements: list[RawElement] = []

        for slide_index, slide in enumerate(presentation.slides):
            slide_num = slide_index + 1
            try:
                slide_elements = self._extract_slide(slide, slide_num)
                elements.extend(slide_elements)
            except Exception:
                logger.exception(
                    "Failed to extract slide",
                    extra={
                        "slide_num": slide_num,
                        "file_key": file_metadata.file_key,
                    },
                )
                continue

        return elements

    def _extract_slide(
        self, slide, slide_num: int
    ) -> list[RawElement]:
        """Extract all elements from a single slide.

        Args:
            slide: A python-pptx slide object.
            slide_num: 1-based slide number.

        Returns:
            List of RawElement from this slide.
        """
        elements: list[RawElement] = []
        position_index = 0

        for shape in slide.shapes:
            try:
                extracted = self._extract_shape(shape, slide_num, position_index)
                if extracted:
                    elements.extend(extracted)
                    position_index += len(extracted)
            except Exception:
                logger.exception(
                    "Failed to extract shape from slide",
                    extra={
                        "slide_num": slide_num,
                        "shape_name": getattr(shape, "name", "unknown"),
                    },
                )
                continue

        return elements

    def _extract_shape(
        self, shape, slide_num: int, position_index: int
    ) -> list[RawElement]:
        """Extract elements from a single shape.

        Args:
            shape: A python-pptx shape object.
            slide_num: 1-based slide number.
            position_index: Current position index within the slide.

        Returns:
            List of RawElement extracted from the shape.
        """
        elements: list[RawElement] = []

        # Table shape
        if shape.has_table:
            table_data = self._extract_table(shape.table)
            if table_data:
                elements.append(
                    RawElement(
                        content=table_data,
                        element_type=ElementType.TABLE,
                        provenance=Provenance(
                            slide_num=slide_num,
                            position_index=position_index + len(elements),
                        ),
                        raw_metadata={"source": "pptx_table"},
                    )
                )
            return elements

        # Image shape
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image_blob = shape.image.blob
                if image_blob:
                    elements.append(
                        RawElement(
                            content=image_blob,
                            element_type=ElementType.IMAGE,
                            provenance=Provenance(
                                slide_num=slide_num,
                                position_index=position_index + len(elements),
                            ),
                            raw_metadata={
                                "source": "pptx_image",
                                "content_type": shape.image.content_type,
                            },
                        )
                    )
            except Exception:
                logger.exception(
                    "Failed to extract image from shape",
                    extra={"slide_num": slide_num},
                )
            return elements

        # Text frame
        if shape.has_text_frame:
            text = self._extract_text_frame(shape.text_frame)
            if text.strip():
                elements.append(
                    RawElement(
                        content=text,
                        element_type=ElementType.TEXT,
                        provenance=Provenance(
                            slide_num=slide_num,
                            position_index=position_index + len(elements),
                        ),
                        raw_metadata={"source": "pptx_text_frame"},
                    )
                )

        return elements

    def _extract_text_frame(self, text_frame) -> str:
        """Extract concatenated text from a text frame.

        Args:
            text_frame: A python-pptx TextFrame object.

        Returns:
            Concatenated text from all paragraphs in the frame.
        """
        paragraphs = []
        for paragraph in text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                paragraphs.append(text)
        return "\n".join(paragraphs)

    def _extract_table(self, table) -> str:
        """Extract table content as a formatted string.

        Args:
            table: A python-pptx Table object.

        Returns:
            Table content formatted as tab-separated rows with newlines,
            or empty string if table is empty.
        """
        rows: list[str] = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cells.append(cell.text.strip())
            rows.append("\t".join(cells))

        return "\n".join(rows) if rows else ""
